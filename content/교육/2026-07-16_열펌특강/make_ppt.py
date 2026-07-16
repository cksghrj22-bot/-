# 열펌의 모든 것 (7/16 특강) — 앳나운 교육개편안 디자인 시스템 정본 이식
# 개편안 실물 기준: 순백 배경 · 검정 볼드 · 골드 점 + 검정 원 강조 · 연회색 카드 ·
# 제목+얇은 구분선 · 하단 중앙 AT NOWN 푸터 · 속 찬 검은 졸라맨 · 큰 중앙 문장 · 노드 플로우.
# 설명글 최소(강사가 말로). 사용: python3 content/교육/2026-07-16_열펌특강/make_ppt.py
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── 개편안 팔레트 ──
BG    = RGBColor(0xFF, 0xFF, 0xFF)   # 순백
INK   = RGBColor(0x1C, 0x1C, 0x1C)   # 검정
SUB   = RGBColor(0x80, 0x80, 0x80)   # 회색 (보조·푸터)
CARD  = RGBColor(0xF2, 0xF1, 0xED)   # 연회색 카드
GOLD  = RGBColor(0xB8, 0x91, 0x2E)   # 골드 점 강조
DIV   = RGBColor(0xE4, 0xE3, 0xDF)   # 구분선
FONT = "Noto Sans CJK KR"

prs = Presentation()
prs.slide_width, prs.slide_height = I(13.333), I(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5


def slide(notes=""):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background(); r.shadow.inherit = False
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def txt(s, x, y, w, h, text, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, spacing=1.0, spc=0.0):
    b = s.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.name = FONT; f.size = Pt(size); f.color.rgb = color; f.bold = bold
        rPr = r._r.get_or_add_rPr()
        if spc:
            rPr.set('spc', str(int(spc * 100)))
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
        ea.set('typeface', FONT)
    return b


def line(s, x1, y1, x2, y2, color=INK, w=2.0):
    ln = s.shapes.add_connector(1, I(x1), I(y1), I(x2), I(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(w); ln.shadow.inherit = False
    return ln


def dline(s, x1, y1, x2, y2, color=GOLD, w=1.6):
    ln = line(s, x1, y1, x2, y2, color=color, w=w)
    etree.SubElement(ln.line._get_or_add_ln(), qn('a:prstDash')).set('val', 'dash')
    return ln


def disc(s, cx, cy, r, fill=INK, line_c=None, w=2.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, I(cx - r), I(cy - r), I(2*r), I(2*r))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_c is None: shp.line.fill.background()
    else: shp.line.color.rgb = line_c; shp.line.width = Pt(w)
    shp.shadow.inherit = False
    return shp


def rrect(s, x, y, w, h, fill=CARD, line_c=None, line_w=1.0, rad=0.06):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    try: shp.adjustments[0] = rad
    except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_c is None: shp.line.fill.background()
    else: shp.line.color.rgb = line_c; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def poly(s, pts, fill=None, line_c=INK, w=2.0, close=True):
    x0, y0 = pts[0]
    fb = s.shapes.build_freeform(Emu(int(I(x0))), Emu(int(I(y0))), scale=1)
    fb.add_line_segments([(Emu(int(I(x))), Emu(int(I(y)))) for x, y in pts[1:]], close=close)
    shp = fb.convert_to_shape()
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_c is None: shp.line.fill.background()
    else: shp.line.color.rgb = line_c; shp.line.width = Pt(w)
    shp.shadow.inherit = False
    return shp


def figure(s, cx, ground_y, scale=1.0):
    """개편안 졸라맨 — 속이 꽉 찬 검은 실루엣이 바닥선 위에 선다."""
    hr = 0.30 * scale
    body_w = 0.52 * scale; body_h = 0.82 * scale
    leg_h = 0.55 * scale
    body_bot = ground_y - leg_h
    body_top = body_bot - body_h
    head_cy = body_top - hr + 0.06 * scale
    # 바닥선
    line(s, cx - 3.4*scale, ground_y, cx + 3.4*scale, ground_y, color=INK, w=1.4)
    # 다리
    line(s, cx - 0.13*scale, body_bot, cx - 0.13*scale, ground_y, color=INK, w=int(7*scale))
    line(s, cx + 0.13*scale, body_bot, cx + 0.13*scale, ground_y, color=INK, w=int(7*scale))
    # 몸통
    rrect(s, cx - body_w/2, body_top, body_w, body_h, fill=INK, rad=0.22)
    # 팔 (살짝 벌림)
    ay = body_top + 0.18*scale
    line(s, cx - body_w/2, ay, cx - body_w/2 - 0.42*scale, ay - 0.02*scale, color=INK, w=int(7*scale))
    line(s, cx + body_w/2, ay, cx + body_w/2 + 0.42*scale, ay - 0.02*scale, color=INK, w=int(7*scale))
    # 머리
    disc(s, cx, head_cy, hr, fill=INK)


def title(s, text, size=28):
    txt(s, 0.9, 0.66, 11.5, 0.9, text, size=size, color=INK, bold=True)
    line(s, 0.9, 1.58, 12.43, 1.58, color=DIV, w=1.4)


def subline(s, text):
    txt(s, 0.9, 1.7, 11.5, 0.5, text, size=13, color=SUB)


def footer(s):
    n = len(prs.slides._sldIdLst)
    txt(s, 0, H - 0.62, W, 0.35, "AT NOWN", size=11, color=SUB, bold=True, align=PP_ALIGN.CENTER, spc=2.0)
    txt(s, W - 1.3, H - 0.62, 0.9, 0.3, f"{n}", size=11, color=SUB, align=PP_ALIGN.RIGHT)


def numbered(s, n, y, head, sub):
    disc(s, 1.25, y + 0.26, 0.25, fill=INK)
    txt(s, 1.0, y + 0.04, 0.5, 0.45, str(n), size=15, color=BG, bold=True, align=PP_ALIGN.CENTER)
    txt(s, 1.85, y, 10, 0.5, head, size=20, color=INK, bold=True)
    txt(s, 1.85, y + 0.48, 10, 0.45, sub, size=13, color=SUB)


def dotcard(s, x, y, w, h, label, sub=""):
    rrect(s, x, y, w, h, fill=CARD)
    disc(s, x + w/2, y + 0.52, 0.06, fill=GOLD)
    txt(s, x, y + h/2 - 0.26, w, 0.5, label, size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        txt(s, x, y + h/2 + 0.2, w, 0.4, sub, size=12, color=SUB, align=PP_ALIGN.CENTER)


def bigcenter(s, text, size=42):
    txt(s, 0.6, 0, 12.1, H - 0.4, text, size=size, color=INK, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)


def node(s, cx, cy, r, label, sub="", dark=False):
    disc(s, cx, cy, r, fill=(INK if dark else CARD), line_c=(None if dark else DIV), w=1.5)
    col = BG if dark else INK
    if sub:
        txt(s, cx - r, cy - 0.32, 2*r, 0.4, label, size=15, color=col, bold=True, align=PP_ALIGN.CENTER)
        txt(s, cx - r, cy + 0.06, 2*r, 0.35, sub, size=12, color=(BG if dark else SUB), align=PP_ALIGN.CENTER)
    else:
        txt(s, cx - r, cy - 0.2, 2*r, 0.45, label, size=16, color=col, bold=True, align=PP_ALIGN.CENTER)


# ══ 1. 갈라짐 (단독) — 검은 프리즘이 하나를 여럿으로 ══
s = slide("오프닝. 갈라짐 도식만. 하나가 여럿으로. 강사가 '오늘 이걸 쪼갭니다'로 연다.")
cx, cy, sz = 6.55, 3.75, 1.4
poly(s, [(cx, cy - 0.66*sz), (cx - 0.58*sz, cy + 0.66*sz), (cx + 0.58*sz, cy + 0.66*sz)],
     fill=None, line_c=INK, w=2.4, close=True)
line(s, cx - 1.9*sz, cy, cx - 0.5*sz, cy, color=INK, w=2.2)
txt(s, cx - 2.5*sz, cy - 0.42, 0.9, 0.7, "?", size=40, color=INK, bold=True, align=PP_ALIGN.CENTER)
for i in range(4):
    t = (i - 1.5) / 1.5
    c = GOLD if i == 0 else INK
    line(s, cx + 0.5*sz, cy, cx + 2.3*sz, cy + t*1.7*sz, color=c, w=2.2)
txt(s, 0, H - 0.62, W, 0.35, "AT NOWN", size=11, color=SUB, bold=True, align=PP_ALIGN.CENTER, spc=2.0)

# ══ 2. 표지 ══
s = slide("제목. 강사가 '오늘 이걸 쪼갭니다'로 연결.")
txt(s, 0.9, 2.55, 11.5, 1.5, "열펌의 모든 것", size=56, color=INK, bold=True)
txt(s, 0.92, 4.05, 8, 0.7, "분해하다", size=23, color=SUB)
txt(s, 0.92, 6.35, 8, 0.4, "차노 · 2026. 7. 16 · 앳나운플레이스 3F", size=14, color=SUB)
txt(s, 0, H - 0.62, W, 0.35, "AT NOWN", size=11, color=SUB, bold=True, align=PP_ALIGN.CENTER, spc=2.0)

# ══ 3. 오늘 (번호 리스트 + 하단 질문 박스) ══
s = slide("오늘은 하나를 깊게 파는 날이 아니라, 다섯을 쪼개보는 날. 강사가 훑는다.")
title(s, "오늘, 무엇을 쪼갤까?")
rows = [("분해", "약과 '열펌'이라는 말부터"),
        ("모질 · 약제", "크림 · 에멀전 · 겔 · 물"),
        ("커트 = 공간", "펌이 살아났다, 그리고 접기"),
        ("Q & A", "궁금한 것부터 함께"),
        ("실험", "모다발로 제품을 비교한다")]
for i, (hd, sb) in enumerate(rows):
    numbered(s, i + 1, 1.95 + i * 0.8, hd, sb)
rrect(s, 0.9, H - 1.35, 11.53, 0.66, fill=CARD, rad=0.12)
txt(s, 0.9, H - 1.2, 11.53, 0.42, "쪼개봐야, 내 것이 된다", size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# ══ 4. '열펌' 글자 쪼개기 ══
s = slide("'열펌'이라는 말부터 쪼갠다. 열|펌. 강사가 열은 왜/펌이란 무엇/롤러볼vs판을 말로.")
title(s, "'열펌'이라는 말부터")
txt(s, 2.7, 2.7, 3.0, 2.0, "열", size=100, color=INK, bold=True, align=PP_ALIGN.CENTER)
dline(s, 6.2, 2.7, 6.2, 5.5, color=GOLD, w=2.0)
txt(s, 6.7, 2.7, 3.0, 2.0, "펌", size=100, color=INK, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# ══ 5. 열 (게이지) ══
s = slide("열마다 하는 일이 다르다. 40-55 / 100+ / 140. 강사가 각 구간 특성 말로.")
title(s, "열마다, 하는 일이 다르다")
by = 4.2
line(s, 1.4, by, 11.9, by, color=DIV, w=1.8)
for off, temp, c in [(0.03, "40–55°C", INK), (0.42, "100°C +", INK), (0.80, "140°C ~", GOLD)]:
    x = 1.4 + off * 10.5
    disc(s, x, by, 0.06, fill=c)
    txt(s, x - 1.1, by - 0.72, 2.2, 0.5, temp, size=19, color=INK, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# ══ 6. 약 (골드점 카드) ══
s = slide("약도 쪼갠다. 연화제·팩·환원제. 강사가 성희 약 분해 해석을 말로.")
title(s, "약도, 쪼개진다")
for i, t in enumerate(["연화제", "팩", "환원제"]):
    dotcard(s, 1.55 + i * 3.5, 3.0, 3.0, 2.0, t)
footer(s)

# ══ 7. 모질·약제 (골드점 카드) ══
s = slide("손상·모질에 약제를 맞춘다. 크림·에멀전·겔·물. 강사가 각 제형 장점 말로.")
title(s, "모질에, 약제를 맞춘다")
for i, t in enumerate(["크림", "에멀전", "겔", "물"]):
    dotcard(s, 0.95 + i * 2.9, 3.0, 2.6, 2.0, t)
footer(s)

# ══ 8. 그림으로 (사진 자리) ══
s = slide("말이 아니라 그림으로. 곱슬 모양표 + 연화·볼륨 변화표. 핀터레스트 이미지로 채운다.")
title(s, "말이 아니라, 그림으로")
for x, lab in [(0.95, "곱슬 모양표"), (6.9, "연화 · 볼륨 변화표")]:
    rrect(s, x, 2.5, 5.5, 3.7, fill=CARD)
    txt(s, x, 4.1, 5.5, 0.5, lab, size=16, color=SUB, align=PP_ALIGN.CENTER)
footer(s)

# ══ 9. 커트=공간 (큰 문장) ══
s = slide("접힐 자리를 만들면 힘이 절반. 강사가 열 100 vs 50, 손상 절반 말로.")
bigcenter(s, "접으면, 힘이 절반", size=46)
footer(s)

# ══ 10. 작용 | 반작용 ══
s = slide("모든 가위질엔 작용과 반작용. 강사가 걷어내면 접힌다 / 너무 걷어내면 부스스 말로.")
title(s, "작용과, 반작용")
txt(s, 1.4, 3.35, 5.0, 1.2, "작용", size=38, color=INK, bold=True, align=PP_ALIGN.CENTER)
txt(s, 6.0, 3.48, 1.3, 0.9, "↔", size=30, color=SUB, align=PP_ALIGN.CENTER)
txt(s, 6.9, 3.35, 5.0, 1.2, "반작용", size=38, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
footer(s)

# ══ 11. 이제야 보인다 (큰 문장) ══
s = slide("'펌이 살아났다'가 아니라 '이제야 보인다'. 공간이 있어야 컬이 보인다. 강사가 말로.")
bigcenter(s, "살아난 게 아니라,\n이제야 보이는 것", size=40)
footer(s)

# ══ 12. Q&A (큰 문장) ══
s = slide("Q&A. 워드월 질문을 함께 푼다.")
bigcenter(s, "궁금한 것부터", size=46)
footer(s)

# ══ 13. 실험 (골드점 카드 — 제품) ══
s = slide("실험. 모다발로 제품 비교. PB·이찌마루·아모스S·시세이도M. 조건은 하나만 바꾼다.")
title(s, "제품을 바꿔, 눈으로 확인한다")
subline(s, "같은 모다발을 나눠, 조건은 하나만 — 그래야 차이가 제품의 것이 된다")
for i, t in enumerate(["PB", "이찌마루", "아모스 S", "시세이도 M"]):
    dotcard(s, 0.95 + i * 2.9, 3.2, 2.6, 2.0, t, "모다발 ½")
footer(s)

# ══ 14. 기록 (노드 플로우) ══
s = slide("한 모다발 = 한 데이터. 진단→선택→결과→다음. 강사가 말로.")
title(s, "한 모다발 = 한 데이터")
ncy, nr = 4.1, 0.72
labels = [("진단", "모질·손상"), ("선택", "제품·온도"), ("결과", "비교 사진"), ("다음", "바꿀 것")]
xs = [2.2, 5.0, 7.8, 10.6]
for i, (lb, sb) in enumerate(labels):
    node(s, xs[i], ncy, nr, lb, sb, dark=(i == 3))
    if i < 3:
        line(s, xs[i] + nr, ncy, xs[i+1] - nr, ncy, color=INK, w=1.6)
footer(s)

# ══ 15. 마무리 (졸라맨 + 큰 문장) ══
s = slide("마무리. 쪼갤수록 쉬워진다. 개편안처럼 홀로 선 아이로 닫는다. 워드월로 한 번 더.")
figure(s, 6.66, 4.7, scale=0.9)
txt(s, 0.6, 5.1, 12.1, 1.2, "쪼갤수록, 쉬워진다", size=35, color=INK, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0, H - 0.62, W, 0.35, "AT NOWN", size=11, color=SUB, bold=True, align=PP_ALIGN.CENTER, spc=2.0)

out = "content/교육/2026-07-16_열펌특강/열펌의모든것_v6.pptx"
prs.save(out)
print(f"저장: {out} · {len(prs.slides._sldIdLst)}장")
